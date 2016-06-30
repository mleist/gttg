# -*- coding: utf-8 -*-
# pylint: disable=C0330,C0301,W0612

from io import BytesIO
from pptx import Presentation
from decouple import config
from django import forms
from django.http import HttpResponse
from django.utils import timezone, formats
from django.contrib import admin
from django.db import models
from django.utils.translation import ugettext as _
from django_object_actions import DjangoObjectActions
from base.models import Gttg
from base import admin_intro, admin_audience, admin_topic
from base.admin_tools import MyHTMLParser

CFG_GTTG_POWERPOINT_TEMPLATE = config(
    'GTTG_POWERPOINT_TEMPLATE',
    default='GTTG Administration',
    cast=str)
CFG_GTTG_PPTX_LAYOUT_TITLE = config('GTTG_PPTX_LAYOUT_TITLE', cast=str)
CFG_GTTG_PPTX_LAYOUT_CONTENT = config('GTTG_PPTX_LAYOUT_CONTENT', cast=str)
CFG_GTTG_PPTX_LAYOUT_LASTPAGE = config('GTTG_PPTX_LAYOUT_LASTPAGE', cast=str)


class GttgAdmin(DjangoObjectActions, admin.ModelAdmin):
    list_display = ('title', 'start_time')
    ordering = ('start_time',)

    formfield_overrides = {
        models.CharField: {'widget': forms.TextInput(attrs={'size': '80'})},
    }

    inlines = [admin_intro.IntroInlineAdmin,
               admin_topic.TopicInlineAdmin,
               admin_audience.AudienceInlineAdmin,
               ]

    def make_ppt_title_page(self, prs, gttg_obj):
        layout = [l for l in prs.slide_layouts if l.name ==
                  CFG_GTTG_PPTX_LAYOUT_TITLE][0]
        # Title
        slide = prs.slides.add_slide(layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = gttg_obj.title.strip()

        subtitle = [ph for ph in shapes.placeholders if ph.name ==
                    'Text Placeholder 1'][0]
        subtitle.text = gttg_obj.subtitle.strip()

    def make_ppt_toc(self, prs, gttg_obj):
        layout = [l for l in prs.slide_layouts if l.name ==
                  CFG_GTTG_PPTX_LAYOUT_CONTENT][0]
        slide = prs.slides.add_slide(layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = gttg_obj.title + ' -\n' + _('overview')
        text_frame = body_shape.text_frame
        para = text_frame.add_paragraph()
        loc_time = timezone.localtime(gttg_obj.start_time)
        para.text = _("Date") + ": " + str(formats.localize(loc_time))
        para.text += " (" + str(loc_time.tzname()) + ")"
        for intro in gttg_obj.intro_set.all():
            para = text_frame.add_paragraph()
            para.text = _("Intro: %s") % intro.title.strip()
        for topic in gttg_obj.topic_set.all():
            para = text_frame.add_paragraph()
            para.text = _("Topic: %s") % topic.title.strip()
        for audience in gttg_obj.audience_set.all():
            para = text_frame.add_paragraph()
            para.text = _("Audience: %s") % audience.title.strip()

    def make_ppt_intro(self, prs, gttg_obj):
        layout = [l for l in prs.slide_layouts if l.name ==
                  CFG_GTTG_PPTX_LAYOUT_CONTENT][0]
        for intro in gttg_obj.intro_set.all():
            slide = prs.slides.add_slide(layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            title_shape.text = gttg_obj.title + ' -\n' + intro.title
            text_frame = body_shape.text_frame
            parser = MyHTMLParser()
            parser.pptx_txBox = text_frame
            parser.feed(intro.content)

    def make_ppt_topics(self, prs, gttg_obj):
        layout = [l for l in prs.slide_layouts if l.name ==
                  CFG_GTTG_PPTX_LAYOUT_CONTENT][0]
        for topic in gttg_obj.topic_set.all():
            slide = prs.slides.add_slide(layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            title_shape.text = gttg_obj.title + ' -\n' + \
                _("Topic: %s") % topic.title.strip()
            text_frame = body_shape.text_frame
            parser = MyHTMLParser()
            parser.pptx_txBox = text_frame
            parser.feed(topic.content)

    def make_ppt_audience(self, prs, gttg_obj):
        layout = [l for l in prs.slide_layouts if l.name ==
                  CFG_GTTG_PPTX_LAYOUT_CONTENT][0]
        for audience in gttg_obj.audience_set.all():
            slide = prs.slides.add_slide(layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            title_shape.text = gttg_obj.title + ' -\n' + \
                _("Audience: %s") % audience.title.strip()
            text_frame = body_shape.text_frame
            parser = MyHTMLParser()
            parser.pptx_txBox = text_frame
            parser.feed(audience.content)

    def make_ppt_last_page(self, prs, gttg_obj):
        layout = [l for l in prs.slide_layouts if l.name ==
                  CFG_GTTG_PPTX_LAYOUT_LASTPAGE][0]
        slide = prs.slides.add_slide(layout)

    def get_ppt(self, request, objects):
        user = None
        if request.user.is_authenticated():
            user = request.user
        else:
            return
        if objects and len(objects) > 0:
            gttg_obj = objects[0]
            prs = Presentation(CFG_GTTG_POWERPOINT_TEMPLATE)
            self.make_ppt_title_page(prs, gttg_obj)
            self.make_ppt_toc(prs, gttg_obj)
            self.make_ppt_intro(prs, gttg_obj)
            self.make_ppt_topics(prs, gttg_obj)
            self.make_ppt_audience(prs, gttg_obj)
            self.make_ppt_last_page(prs, gttg_obj)
            output = BytesIO()
            prs.save(output)
            response = HttpResponse(output.getvalue(),
                                    content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
            response['Content-Disposition'] = 'attachment; filename="gttg.pptx"'
            response['Content-Length'] = output.tell()
            return response

    get_ppt.label = _("Get PPT")
    get_ppt.short_description = _("Get as PowerPoint")
    actions = ('get_ppt', )


admin.site.register(Gttg, GttgAdmin)
